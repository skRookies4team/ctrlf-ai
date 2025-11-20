"""
파일 파서 - 다양한 형식 지원 (PDF, HWP, DOCX, PPTX)

통합 히스토리:
- 기존: pypdf 기반 PDF 전용
- 업데이트: pdfplumber 기반 PDF + langflow_세희 코드에서 HWP 파서 통합
- HWP 어댑터: core/hwp_converter.py로 여러 변환 방법 통합
"""
import logging
import os
from pathlib import Path
from typing import Optional

# HWP 변환 어댑터 import
from core.hwp_converter import convert_hwp_to_text, get_available_methods

logger = logging.getLogger(__name__)

# ========================================
# 선택적 의존성 체크
# ========================================

# pdfplumber (PDF 파싱)
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False
    logger.warning("pdfplumber not installed. PDF parsing will use fallback method.")

# HWP 변환은 hwp_converter.py에서 자동 선택
# (hwp5txt → LibreOffice → pyhwp 순으로 시도)
hwp_methods = get_available_methods()
if hwp_methods:
    logger.info(f"HWP converters available: {hwp_methods}")
else:
    logger.warning("No HWP converter available. HWP files will return empty text.")

# python-docx (DOCX 파싱)
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not installed. DOCX files will be skipped.")

# python-pptx (PPTX 파싱)
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not installed. PPTX files will be skipped.")


# ========================================
# PDF 파서 (pdfplumber 기반)
# ========================================

def extract_text_from_pdf(pdf_path: str) -> str:
    """
    PDF 파일에서 텍스트 추출 (pdfplumber 사용)

    세희 파서에서 가져온 부분:
    - pdfplumber 기반 추출 로직
    - 페이지별 순회 및 텍스트 결합

    Args:
        pdf_path: PDF 파일 경로

    Returns:
        str: 추출된 텍스트

    Raises:
        RuntimeError: PDF 파싱 실패 시
    """
    try:
        path = Path(pdf_path)

        # 파일 존재 여부 확인
        if not path.exists():
            raise RuntimeError(f"File not found: {pdf_path}")

        logger.info(f"Parsing PDF: {pdf_path}")

        # pdfplumber 사용 (세희 코드 기반)
        if PDFPLUMBER_AVAILABLE:
            text_parts = []
            with pdfplumber.open(pdf_path) as pdf:
                page_count = len(pdf.pages)
                logger.info(f"PDF has {page_count} pages")

                for page_num, page in enumerate(pdf.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(page_text)
                        else:
                            text_parts.append("")
                    except Exception as e:
                        logger.warning(f"Failed to extract text from page {page_num}: {e}")
                        text_parts.append("")

            full_text = "\n".join(text_parts)
            logger.info(f"Successfully extracted {len(full_text)} characters from PDF")
            return full_text

        # Fallback: pypdf
        else:
            from pypdf import PdfReader
            logger.info("Using fallback pypdf parser")

            reader = PdfReader(str(path))
            page_count = len(reader.pages)
            logger.info(f"PDF has {page_count} pages")

            text_parts = []
            for page_num, page in enumerate(reader.pages):
                try:
                    text = page.extract_text()
                    if text is None:
                        text = ""
                    text_parts.append(text)
                except Exception as e:
                    logger.warning(f"Failed to extract text from page {page_num}: {e}")
                    text_parts.append("")

            full_text = "\n".join(text_parts)
            logger.info(f"Successfully extracted {len(full_text)} characters from PDF")
            return full_text

    except Exception as e:
        error_msg = f"Failed to parse PDF {pdf_path}: {str(e)}"
        logger.error(error_msg)
        raise RuntimeError(error_msg)


# ========================================
# HWP 파서 (pyhwp 기반)
# ========================================

def extract_text_from_hwp(hwp_path: str) -> str:
    """
    HWP 파일에서 텍스트 추출

    ⚠️ core/hwp_converter.py 어댑터 사용 (세희 코드 기반)
    ⚠️ 여러 변환 방법 중 자동 선택:
       1. hwp5txt (세희 방식, Linux/Docker 권장)
       2. LibreOffice CLI (크로스 플랫폼)
       3. pyhwp (Python 2, deprecated)

    Args:
        hwp_path: HWP 파일 경로

    Returns:
        str: 추출된 텍스트 (실패 시 빈 문자열)

    Note:
        변환 방법이 없으면 경고 로그를 출력하고 빈 문자열 반환
    """
    try:
        path = Path(hwp_path)

        # 파일 존재 여부 확인
        if not path.exists():
            logger.error(f"File not found: {hwp_path}")
            return ""

        logger.info(f"Parsing HWP: {hwp_path}")

        # hwp_converter 어댑터 호출 (자동 변환 방법 선택)
        text = convert_hwp_to_text(str(path))

        if text:
            logger.info(f"Successfully extracted {len(text)} characters from HWP")
        else:
            logger.warning(f"No text extracted from HWP: {hwp_path}")

        return text

    except Exception as e:
        error_msg = f"Failed to parse HWP {hwp_path}: {str(e)}"
        logger.error(error_msg)
        # HWP 파싱 실패는 치명적이지 않으므로 빈 문자열 반환
        return ""


# ========================================
# DOCX 파서 (python-docx 기반)
# ========================================

def extract_text_from_docx(docx_path: str) -> str:
    """
    DOCX 파일에서 텍스트 추출 (python-docx 사용)

    Args:
        docx_path: DOCX 파일 경로

    Returns:
        str: 추출된 텍스트 (python-docx 없으면 빈 문자열)
    """
    if not DOCX_AVAILABLE:
        logger.warning(f"python-docx not installed. Skipping DOCX file: {docx_path}")
        return ""

    try:
        path = Path(docx_path)

        if not path.exists():
            logger.error(f"File not found: {docx_path}")
            return ""

        logger.info(f"Parsing DOCX: {docx_path}")

        doc = Document(str(path))
        text_parts = []

        # 문단 추출
        for para in doc.paragraphs:
            text_parts.append(para.text)

        # 표(table) 추출
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text_parts.append(cell.text)

        full_text = "\n".join(text_parts)
        logger.info(f"Successfully extracted {len(full_text)} characters from DOCX")
        return full_text

    except Exception as e:
        error_msg = f"Failed to parse DOCX {docx_path}: {str(e)}"
        logger.error(error_msg)
        return ""


# ========================================
# PPTX 파서 (python-pptx 기반)
# ========================================

def extract_text_from_pptx(pptx_path: str) -> str:
    """
    PPTX 파일에서 텍스트 추출 (python-pptx 사용)

    Args:
        pptx_path: PPTX 파일 경로

    Returns:
        str: 추출된 텍스트 (python-pptx 없으면 빈 문자열)
    """
    if not PPTX_AVAILABLE:
        logger.warning(f"python-pptx not installed. Skipping PPTX file: {pptx_path}")
        return ""

    try:
        path = Path(pptx_path)

        if not path.exists():
            logger.error(f"File not found: {pptx_path}")
            return ""

        logger.info(f"Parsing PPTX: {pptx_path}")

        prs = Presentation(str(path))
        text_parts = []

        # 각 슬라이드 순회
        for slide_num, slide in enumerate(prs.slides):
            # 슬라이드 내 모든 shape 확인
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)

        full_text = "\n".join(text_parts)
        logger.info(f"Successfully extracted {len(full_text)} characters from PPTX")
        return full_text

    except Exception as e:
        error_msg = f"Failed to parse PPTX {pptx_path}: {str(e)}"
        logger.error(error_msg)
        return ""


# ========================================
# 통합 파서 (확장자 기반 라우팅)
# ========================================

def extract_text_from_file(file_path: str) -> str:
    """
    파일 형식에 따라 적절한 파서를 선택하여 텍스트 추출

    지원 형식:
    - .pdf: pdfplumber 기반
    - .hwp: pyhwp 기반 (세희 파서)
    - .docx: python-docx 기반
    - .pptx: python-pptx 기반

    Args:
        file_path: 파일 경로

    Returns:
        str: 추출된 텍스트

    Raises:
        ValueError: 지원하지 않는 파일 형식
        RuntimeError: 파싱 실패 (PDF의 경우)
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    logger.info(f"Detecting file format: {ext}")

    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.hwp':
        return extract_text_from_hwp(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    else:
        error_msg = f"Unsupported file format: {ext}"
        logger.error(error_msg)
        raise ValueError(error_msg)


# ========================================
# 페이지 수 추출 헬퍼 (PDF 전용)
# ========================================

def get_page_count(pdf_path: str) -> Optional[int]:
    """
    PDF 파일의 페이지 수 반환

    Args:
        pdf_path: PDF 파일 경로

    Returns:
        Optional[int]: 페이지 수 (실패 시 None)
    """
    try:
        if PDFPLUMBER_AVAILABLE:
            with pdfplumber.open(pdf_path) as pdf:
                return len(pdf.pages)
        else:
            from pypdf import PdfReader
            reader = PdfReader(pdf_path)
            return len(reader.pages)
    except Exception as e:
        logger.error(f"Failed to get page count from {pdf_path}: {e}")
        return None
